import os
from pptx import Presentation
from pptx.util import Inches

def add_logos_to_presentation(input_path, logo_path, output_path, logo_width_inches=0.6, margin_inches=0.15):
    if not os.path.exists(input_path):
        print(f"Error: {input_path} not found.")
        return False
    if not os.path.exists(logo_path):
        print(f"Error: {logo_path} not found.")
        return False

    prs = Presentation(input_path)
    logo_width = Inches(logo_width_inches)
    margin = Inches(margin_inches)

    for i, slide in enumerate(prs.slides):
        # Top-left logo
        slide.shapes.add_picture(
            logo_path,
            margin,
            margin,
            width=logo_width
        )
        
        # Top-right logo
        left_pos = prs.slide_width - logo_width - margin
        slide.shapes.add_picture(
            logo_path,
            left_pos,
            margin,
            width=logo_width
        )
        print(f"  Added logos to slide {i+1}")

    prs.save(output_path)
    print(f"Saved: {output_path}")
    return True

if __name__ == "__main__":
    logo = "../afit logo.jpg"
    ppt_file = "../EORIS_Defense.pptx"
    output_ppt = "../EORIS_Defense_with_logos.pptx"
    
    print(f"Adding logos to {ppt_file}...")
    success = add_logos_to_presentation(ppt_file, logo, output_ppt)
    if success:
        print("Logo insertion completed successfully!")
    else:
        print("Failed to add logos.")
